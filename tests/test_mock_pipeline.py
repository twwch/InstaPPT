import os
import json
import pytest
from unittest.mock import MagicMock, patch
from instappt.core import PPTTranslator
from instappt.models import SDKConfig, ModelConfig, EvaluationResult, EvaluationMetrics

@pytest.fixture
def mock_config():
    cfg = ModelConfig(base_url="http://mock", api_key="mock", model="mock")
    return SDKConfig(translator_config=cfg, optimizer_config=cfg, evaluator_config=cfg)

class MockOpenAIClient:
    def __init__(self):
        self.responses = []
        self.call_count = 0

    def add_response(self, content):
        self.responses.append(content)

    def chat_completion(self, config, prompt, tag=None):
        # Return next response in queue
        if self.call_count < len(self.responses):
            resp = self.responses[self.call_count]
            self.call_count += 1
            return resp
        return "Generic Response"

def test_smart_selection_fallback(mock_config):
    """Test that if Optimization score is lower than Translation, we strictly fallback to Translation."""
    translator = PPTTranslator(mock_config)
    mock_client = MockOpenAIClient()
    translator._call_llm = mock_client.chat_completion
    
    # Setup mocks
    # 1. Translation A
    mock_client.add_response("Translation A")
    # 2. Eval A (Score 9.0)
    mock_client.add_response('{"original_text": "src", "translated_text": "A", "metrics": {"accuracy": 9, "fluency": 9, "consistency": 9, "terminology": 9, "completeness": 9}, "overall_score": 9.0, "suggestions": "None"}')
    # 3. Optimization C
    mock_client.add_response("Optimization C-Worse")
    # 4. Eval C (Score 7.0)
    mock_client.add_response('{"original_text": "src", "translated_text": "C", "metrics": {"accuracy": 7, "fluency": 7, "consistency": 7, "terminology": 7, "completeness": 7}, "overall_score": 7.0, "suggestions": "Bad"}')
    
    # Run pipeline for one segment
    items = {"id1": "Source Text"}
    
    # Note: Staged processing runs concurrently, so order of calls might vary if we had multiple items.
    # With 1 item, the order within stage is deterministic, but across stages it is sequential (A -> Eval A -> Opt C -> Eval C).
    
    # Stage A
    seg_map = translator._stage_translation(items, "Chinese")
    # Stage B
    translator._stage_evaluation(seg_map, "Chinese")
    # Stage C
    translator._stage_optimization(seg_map, "Chinese")
    
    res = seg_map["id1"]
    
    # Result must be A because 9.0 > 7.0
    assert res.final_text == "Translation A"
    assert res.optimized_text_c == "Optimization C-Worse"
    assert res.evaluation_a.overall_score == 9.0
    assert res.evaluation_c.overall_score == 7.0

def test_pipeline_end_to_end(mock_config, tmp_path):
    translator = PPTTranslator(mock_config)
    # Mock OpenAI client
    with patch('instappt.core.OpenAI') as MockOpenAI:
        mock_client = MockOpenAI.return_value
        
        def mock_create(*args, **kwargs):
            messages = kwargs.get('messages', [])
            prompt = messages[0]['content'] if messages else ""
            
            content = "Generic Response"
            # More specific matching to avoid overlap
            if "Translate the following" in prompt:
                content = "Translated Text"
            elif "linguistic evaluator" in prompt or "Evaluate the translation" in prompt:
                content = json.dumps({
                    "accuracy": 8, "fluency": 9, "consistency": 8, "terminology": 7, "completeness": 10,
                    "suggestions": "Good job", "overall_score": 8.4
                })
            elif "suggestions" in prompt or "refined translation" in prompt:
                 content = "Optimized Translated Text"
            
            mock_response = MagicMock()
            mock_response.choices[0].message.content = content
            mock_response.usage.prompt_tokens = 10
            mock_response.usage.completion_tokens = 20
            mock_response.usage.total_tokens = 30
            mock_response.usage.model_dump.return_value = {"prompt_tokens": 10, "completion_tokens": 20}
            return mock_response

        mock_client.chat.completions.create.side_effect = mock_create
            
        # Run process (assuming sample.pptx exists in cwd, otherwise create one)
        input_ppt = "sample.pptx"
        output_ppt = str(tmp_path / "output.pptx")
        
        if not os.path.exists(input_ppt):
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Hello World"
            prs.save(input_ppt)
            
        translator.process_ppt(input_ppt, output_ppt, "German")
    
    # Check assertions
    assert len(translator.segments) > 0
    assert translator.segments[0].translated_text_a == "Translated Text"
    assert translator.segments[0].optimized_text_c == "Optimized Translated Text"


def test_pipeline_skip_optimization(mock_config, tmp_path):
    translator = PPTTranslator(mock_config)

    # Mock _call_llm to return high score
    def mock_call_llm(config, prompt, tag=None):
        if "Translate" in prompt:
            return "Perfect Text"
        elif "evaluate" in prompt or "Evaluate" in prompt:
            return json.dumps({
                "accuracy": 10, "fluency": 10, "consistency": 10, "terminology": 10, "completeness": 10,
                "suggestions": "None", "overall_score": 9.8 
            })
        elif "Optimize" in prompt:
             return "Should Not Be Called"
        return "Generic Response"

    translator._call_llm = MagicMock(side_effect=mock_call_llm)

    input_ppt = "sample_perf.pptx"
    output_ppt = str(tmp_path / "output_perf.pptx")

    if not os.path.exists(input_ppt):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Hello World"
        prs.save(input_ppt)

    translator.process_ppt(input_ppt, output_ppt, "German")

    # Optimization should be skipped
    assert len(translator.segments) > 0
    assert translator.segments[0].translated_text_a == "Perfect Text"
    assert translator.segments[0].evaluation_a.overall_score == 9.8
    assert translator.segments[0].optimized_text_c == "" # Should be empty
    assert translator.segments[0].final_text == "Perfect Text"
    
    # Check reports
    output_dir = str(tmp_path / "reports")
    translator.generate_reports(output_dir)
    
    assert os.path.exists(os.path.join(output_dir, "comparison.json"))
    assert os.path.exists(os.path.join(output_dir, "evaluation_report.xlsx"))
    assert os.path.exists(os.path.join(output_dir, "comparison.pdf"))
    assert os.path.exists(os.path.join(output_dir, "token_usage.json"))

if __name__ == "__main__":
    pytest.main([__file__])
