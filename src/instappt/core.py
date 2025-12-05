from typing import List, Dict, Any
from .models import SDKConfig, TranslationSegment, ModelConfig, EvaluationResult, EvaluationMetrics
from .prompts import TRANSLATION_PROMPT, OPTIMIZATION_PROMPT, EVALUATION_PROMPT
import pptx
from openai import OpenAI
import json
import re
import time
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from tqdm import tqdm
from threading import Lock


class PPTTranslator:
    def __init__(self, config: SDKConfig, concurrency: int = 32):
        self.config = config
        self.concurrency = concurrency
        self.segments: List[TranslationSegment] = []
        self.detailed_logs: List[Dict[str, Any]] = [] # Detailed telemetry
        self.log_lock = Lock() # Thread safety for logs
        
        # Caching
        self.cache_file = ".instappt_cache.json"
        self.cache_lock = Lock()
        self.cache = self._load_cache()

    def _load_cache(self) -> Dict[str, str]:
        if os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                print(f"Warning: Failed to load cache: {e}")
        return {}

    def _save_cache(self, key: str, value: Any):
        if not self.config.enable_cache:
            return
            
        with self.cache_lock:
            self.cache[key] = value
            try:
                with open(self.cache_file, 'w', encoding='utf-8') as f:
                    json.dump(self.cache, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"Warning: Failed to save cache: {e}")

    def process_ppt(self, input_path: str, output_path: str, target_language: str, progress_callback=None):
        """
        Main pipeline: Load -> Translate -> Evaluate -> Optimize -> Save
        """
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")
            
        presentation = pptx.Presentation(input_path)
        
        # 1. Extract and map all text
        text_map = self._extract_text_and_map(presentation)
        
        # Filter empty
        items_to_process = {sid: txt for sid, txt in text_map.items() if txt.strip()}
        
        # Initialize segments
        self.segments = []
        # We use a dict to track segments by ID during processing
        segment_map: Dict[str, TranslationSegment] = {}
        
        print(f"Processing {len(items_to_process)} segments with concurrency {self.concurrency}...")

        # --- Stage A: Translation (Green) ---
        segment_map = self._stage_translation(items_to_process, target_language, progress_callback)

        # --- Stage B: Evaluation (Blue) ---
        # We process the segments generated in Stage A
        self._stage_evaluation(segment_map, target_language, progress_callback)

        # --- Stage C: Optimization (Yellow) ---
        self._stage_optimization(segment_map, target_language, progress_callback)

        # Finalize list
        # Ensure order matches items_to_process if possible, or just list
        self.segments = list(segment_map.values())

        # 3. Update text sequentially to avoid race conditions on PPT object
        for sid, segment in segment_map.items():
            if segment.final_text:
                self._replace_text(presentation, sid, segment.final_text)
        
        # 4. Save new PPT
        output_dir = os.path.dirname(output_path)
        if not os.path.exists(output_dir):
            print(f"WARNING: Output directory {output_dir} was missing. Re-creating it.")
            os.makedirs(output_dir)
            
        presentation.save(output_path)

    def _stage_translation(self, items: Dict[str, str], target_lang: str, progress_callback=None) -> Dict[str, TranslationSegment]:
        results = {}
        with ThreadPoolExecutor(max_workers=self.concurrency) as executor:
            future_to_sid = {}
            for sid, text in items.items():
                future = executor.submit(self._task_translate, text, target_lang)
                future_to_sid[future] = sid

            # Green bar
            processed = 0
            total = len(items)
            with tqdm(total=total, desc="Translation", unit="seg", colour='green') as pbar:
                for future in as_completed(future_to_sid):
                    sid = future_to_sid[future]
                    try:
                        seg = future.result()
                        results[sid] = seg
                    except Exception as e:
                        print(f"Error translating {sid}: {e}")
                    
                    pbar.update(1)
                    processed += 1
                    if progress_callback:
                        try:
                            progress_callback(processed, total, "Translation")
                        except Exception as e:
                            print(f"Callback error: {e}")
        return results

    def _task_translate(self, text: str, target_lang: str) -> TranslationSegment:
        start_a = time.time()
        trans_prompt = TRANSLATION_PROMPT.format(
            source_language="Auto",
            target_language=target_lang,
            text=text
        )
        trans_a = self._call_llm(self.config.translator_config, trans_prompt, tag="Stage A: Translation")
        dur_a = time.time() - start_a
        
        # Create initial segment state
        return TranslationSegment(
            original_text=text,
            translated_text_a=trans_a,
            duration_a=dur_a,
            evaluation_a=None, # Filled next
            optimized_text_c=None,
            duration_c=0.0,
            evaluation_c=None,
            final_text=None # Decided later
        )

    def _stage_evaluation(self, segment_map: Dict[str, TranslationSegment], target_lang: str, progress_callback=None):
        with ThreadPoolExecutor(max_workers=self.concurrency) as executor:
            future_to_sid = {
                executor.submit(self._task_evaluate, seg, target_lang): sid 
                for sid, seg in segment_map.items()
            }
            # Blue bar
            processed = 0
            total = len(segment_map)
            with tqdm(total=total, desc="Evaluation ", unit="seg", colour='blue') as pbar:
                for future in as_completed(future_to_sid):
                    # We modify segment in place or return update
                    # Returning update is safer
                    sid = future_to_sid[future]
                    try:
                        updated_seg = future.result()
                        segment_map[sid] = updated_seg
                    except Exception as e:
                         print(f"Error in Evaluation stage for {sid}: {e}")
                    pbar.update(1)
                    processed += 1
                    if progress_callback:
                        progress_callback(processed, total, "Evaluation")

    def _task_evaluate(self, seg: TranslationSegment, target_lang: str) -> TranslationSegment:
        start_eval_a = time.time()
        eval_prompt_a = EVALUATION_PROMPT.format(
            source_text=seg.original_text, 
            target_language=target_lang, 
            translation=seg.translated_text_a
        )
        eval_resp_a = self._call_llm(self.config.evaluator_config, eval_prompt_a, tag="Stage B1: Eval of A")
        eval_result_a = self._parse_evaluation(eval_resp_a)
        dur_eval_a = time.time() - start_eval_a
        if eval_result_a: eval_result_a.duration_seconds = dur_eval_a
        
        seg.evaluation_a = eval_result_a
        return seg

    def _stage_optimization(self, segment_map: Dict[str, TranslationSegment], target_lang: str, progress_callback=None):
        with ThreadPoolExecutor(max_workers=self.concurrency) as executor:
            future_to_sid = {
                executor.submit(self._task_optimize, seg, target_lang): sid 
                for sid, seg in segment_map.items()
            }
            # Yellow bar
            processed = 0
            total = len(segment_map)
            with tqdm(total=total, desc="Optimization", unit="seg", colour='yellow') as pbar:
                for future in as_completed(future_to_sid):
                    sid = future_to_sid[future]
                    try:
                         updated_seg = future.result()
                         segment_map[sid] = updated_seg
                    except Exception as e:
                         print(f"Error in Optimization stage for {sid}: {e}")
                    pbar.update(1)
                    processed += 1
                    if progress_callback:
                        progress_callback(processed, total, "Optimization")

    def _task_optimize(self, seg: TranslationSegment, target_lang: str) -> TranslationSegment:
        # Decision logic
        if not seg.evaluation_a or seg.evaluation_a.overall_score > 9.5:
            # Skip
            seg.final_text = seg.translated_text_a
            seg.optimized_text_c = ""
            return seg # Return early
            
        # Optimize
        start_c = time.time()
        opt_prompt = OPTIMIZATION_PROMPT.format(
            target_language=target_lang,
            source_text=seg.original_text,
            translation=seg.translated_text_a,
            suggestions=seg.evaluation_a.suggestions
        )
        trans_c = self._call_llm(self.config.optimizer_config, opt_prompt, tag="Stage C: Optimization")
        dur_c = time.time() - start_c
        
        # Evaluate C (B2)
        start_eval_c = time.time()
        eval_prompt_c = EVALUATION_PROMPT.format(
            source_text=seg.original_text, 
            target_language=target_lang, 
            translation=trans_c
        )
        eval_resp_c = self._call_llm(self.config.evaluator_config, eval_prompt_c, tag="Stage B2: Eval of C")
        eval_result_c = self._parse_evaluation(eval_resp_c)
        dur_eval_c = time.time() - start_eval_c
        if eval_result_c: eval_result_c.duration_seconds = dur_eval_c
        
        seg.optimized_text_c = trans_c
        seg.duration_c = dur_c
        seg.evaluation_c = eval_result_c
        
        # Smart Selection: Choose the version with the higher score
        score_a = seg.evaluation_a.overall_score if seg.evaluation_a else 0
        score_c = eval_result_c.overall_score if eval_result_c else 0
        
        if score_c >= score_a and trans_c.strip():
            seg.final_text = trans_c
        else:
            # Revert to A if optimization failed to improve score (or was worse)
            seg.final_text = seg.translated_text_a
        
        return seg

    def _extract_text_and_map(self, prs) -> Dict[str, str]:
        """
        Traverses the presentation and builds a map of {unique_id: text}.
        Unique ID structure: slide_idx:shape_id:paragraph_idx
        """
        text_map = {}
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                self._process_shape(shape, slide_idx, text_map)
        return text_map

    def _process_shape(self, shape, slide_idx, text_map):
        if shape.has_text_frame:
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                 # ID format: slide_index:shape_id:paragraph_index
                 # Use runs? For simplicity, we translate paragraph level.
                 full_text = "".join(run.text for run in paragraph.runs)
                 if full_text.strip():
                     segment_id = f"{slide_idx}:{shape.shape_id}:{p_idx}"
                     text_map[segment_id] = full_text
        
        if shape.has_table:
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    # Recursively treat cell as a text container, but it's a bit complex with shape_id.
                    # Cell doesn't have a shape_id unique in the slide easily usable.
                    # We use table_shape_id:row:col:p_idx
                    for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                         full_text = "".join(run.text for run in paragraph.runs)
                         if full_text.strip():
                             segment_id = f"{slide_idx}:{shape.shape_id}:table:{row_idx}:{col_idx}:{p_idx}"
                             text_map[segment_id] = full_text

        if shape.shape_type == 6: # Group
             # Recursive for groups if needed, but python-pptx group support is tricky.
             # Check if .shapes exists
             if hasattr(shape, "shapes"):
                 for child_shape in shape.shapes:
                     self._process_shape(child_shape, slide_idx, text_map)

    def _replace_text(self, prs, segment_id: str, new_text: str):
        """
        Parse ID and replace text.
        """
        parts = segment_id.split(":")
        slide_idx = int(parts[0])
        shape_id = int(parts[1])
        slide = prs.slides[slide_idx]
        
        # Find shape
        target_shape = None
        # This is O(N) per replacement, could be optimized with a pre-calc map, but fine for PPT sizes.
        
        # Helper to find shape recursively
        def find_shape(shapes, s_id):
            for shape in shapes:
                if shape.shape_id == s_id:
                    return shape
                if shape.shape_type == 6 and hasattr(shape, "shapes"):
                    found = find_shape(shape.shapes, s_id)
                    if found: return found
            return None

        target_shape = find_shape(slide.shapes, shape_id)
        if not target_shape:
            return

        if "table" in parts:
             # Table replacement
             row_idx = int(parts[3])
             col_idx = int(parts[4])
             p_idx = int(parts[5])
             cell = target_shape.table.rows[row_idx].cells[col_idx]
             if p_idx < len(cell.text_frame.paragraphs):
                 paragraph = cell.text_frame.paragraphs[p_idx]
                 self._set_paragraph_text(paragraph, new_text, self._get_original_text_for_segment(prs, segment_id))
        else:
             # TextFrame replacement
             p_idx = int(parts[2])
             if p_idx < len(target_shape.text_frame.paragraphs):
                 paragraph = target_shape.text_frame.paragraphs[p_idx]
                 self._set_paragraph_text(paragraph, new_text, self._get_original_text_for_segment(prs, segment_id))

    def _get_original_text_for_segment(self, prs, segment_id: str) -> str:
        """
        Retrieves the original text for a given segment_id.
        This is a temporary helper to get the original text for autosizing.
        In a real scenario, this would be passed down from the segment_map.
        """
        parts = segment_id.split(":")
        slide_idx = int(parts[0])
        shape_id = int(parts[1])
        slide = prs.slides[slide_idx]

        def find_shape(shapes, s_id):
            for shape in shapes:
                if shape.shape_id == s_id:
                    return shape
                if shape.shape_type == 6 and hasattr(shape, "shapes"):
                    found = find_shape(shape.shapes, s_id)
                    if found: return found
            return None

        target_shape = find_shape(slide.shapes, shape_id)
        if not target_shape:
            return ""

        if "table" in parts:
            row_idx = int(parts[3])
            col_idx = int(parts[4])
            p_idx = int(parts[5])
            cell = target_shape.table.rows[row_idx].cells[col_idx]
            if p_idx < len(cell.text_frame.paragraphs):
                return "".join(run.text for run in cell.text_frame.paragraphs[p_idx].runs)
        else:
            p_idx = int(parts[2])
            if p_idx < len(target_shape.text_frame.paragraphs):
                return "".join(run.text for run in target_shape.text_frame.paragraphs[p_idx].runs)
        return ""

    def _set_paragraph_text(self, paragraph, new_text, original_text=""):
        from pptx.util import Pt

        # Preserve first run's formatting if possible, clear others
        if not paragraph.runs:
            paragraph.add_run().text = new_text
            return
            
        # Basic font preservation from first run
        p_font = paragraph.runs[0].font
        font_name = p_font.name
        font_size = p_font.size
        bold = p_font.bold
        italic = p_font.italic
        # Capture color (RGB or Theme)
        color_rgb = None
        color_theme = None
        color_brightness = 0.0
        
        if p_font.color:
            try:
                if p_font.color.type == 1: # MSO_COLOR_TYPE.RGB
                    color_rgb = p_font.color.rgb
                elif p_font.color.type == 2: # MSO_COLOR_TYPE.SCHEME
                    color_theme = p_font.color.theme_color
                    color_brightness = p_font.color.brightness
            except:
                pass

        paragraph.clear() # Removes all runs
        new_run = paragraph.add_run()
        new_run.text = new_text
        
        # Apply style
        # Fix for Chinese font rendering on macOS/LibreOffice
        # If text contains non-ASCII characters, force a compatible font
        is_non_ascii = any(ord(c) > 127 for c in new_text)
        if is_non_ascii:
            # Arial Unicode MS is widely available on macOS and supports Chinese
            # SimHei or Microsoft YaHei might be better if installed, but Arial Unicode MS is safer default
            # Update: User reported issues. Switching to 'Heiti SC' (STHeiti) which is confirmed present.
            new_run.font.name = "Heiti SC" 
        elif font_name: 
            new_run.font.name = font_name
        if font_size: new_run.font.size = font_size
        new_run.font.bold = bold
        new_run.font.italic = italic
        
        # Apply Color
        if color_rgb:
            new_run.font.color.rgb = color_rgb
        elif color_theme:
            new_run.font.color.theme_color = color_theme
            if color_brightness:
                new_run.font.color.brightness = color_brightness
        
        # Font Autosizing
        if original_text and len(new_text) > len(original_text) * 1.1 and font_size:
             ratio = len(new_text) / len(original_text)
             scale = 1.0 / (ratio ** 0.5)
             new_size = font_size * scale
             # Enforce minimum size of 8pt
             new_run.font.size = max(new_size, Pt(8))

    def _call_llm(self, config: ModelConfig, prompt: str, tag: str = "General") -> str:
        import hashlib
        
        # 1. Check Cache
        cache_key = hashlib.md5((config.model + prompt).encode('utf-8')).hexdigest()
        
        with self.cache_lock:
            if self.config.enable_cache and cache_key in self.cache:
                # CACHE HIT
                # We still log it but with 0 duration/cost to indicate hit
                with self.log_lock:
                    self.detailed_logs.append({
                        "timestamp": time.time(), # Added timestamp for consistency
                        "tag": tag,
                        "model": config.model, # Added for compatibility
                        "model_config_alias": config.model, # Use config model as we don't have real model from cache
                        "duration": 0.0,
                        "messages": [{"role": "user", "content": prompt}],
                        "response_content": self.cache[cache_key], # Log the cached content
                        "usage": { # Mimic usage structure for consistency
                            "prompt_tokens": 0,
                            "completion_tokens": 0,
                            "total_tokens": 0,
                            "cached": True
                        }
                    })
                return self.cache[cache_key]

        # 2. Cache Miss - Call API
        messages = [{"role": "user", "content": prompt}]
        # print(f"DEBUG: Calling API with model {config.model}, key_len={len(config.api_key)}")
        client = OpenAI(api_key=config.api_key, base_url=config.base_url)
        
        start_time = time.time()
        try:
            # print("DEBUG: Sending request...")
            response = client.chat.completions.create(
                model=config.model,
                messages=messages,
                temperature=0.3
            )
            duration = time.time() - start_time
            content = response.choices[0].message.content
            
            # Log usage
            real_model = response.model if hasattr(response, 'model') else config.model
            usage = response.usage
            usage_dict = usage.model_dump() if hasattr(usage, 'model_dump') else usage.__dict__
            
            log_entry = {
                "timestamp": time.time(),
                "tag": tag,
                "model": real_model, # Key for reports
                "model_config_alias": config.model,
                "duration": duration,
                "messages": messages,
                "response_content": content,
                "usage": usage_dict
            }
        
            with self.log_lock:
                self.detailed_logs.append(log_entry)
            
            # 3. Update Cache
            # 3. Update Cache
            self._save_cache(cache_key, content)
            
            return content
            
        except Exception as e:
            print(f"LLM Error ({tag}): {e}")
            return ""

    def _parse_evaluation(self, json_str: str) -> EvaluationResult:
        """Robust JSON parsing for evaluation result."""
        try:
            # Try to find JSON block if mixed with text
            match = re.search(r'\{.*\}', json_str, re.DOTALL)
            if match:
                json_str = match.group(0)
            data = json.loads(json_str)
            
            # Local import avoided by top-level import, but kept logic same
            # Handle nested metrics if present (as requested in prompt)
            metrics_data = data.get('metrics', data)
            
            metrics = EvaluationMetrics(
                accuracy=metrics_data.get('accuracy', 0),
                fluency=metrics_data.get('fluency', 0),
                consistency=metrics_data.get('consistency', 0),
                terminology=metrics_data.get('terminology', 0),
                completeness=metrics_data.get('completeness', 0)
            )
            return EvaluationResult(
                metrics=metrics,
                suggestions=data.get('suggestions', 'No suggestions provided'),
                overall_score=float(data.get('overall_score', 0.0))
            )
        except Exception as e:
            # Fallback for failed parsing
            print(f"Failed to parse evaluation JSON: {e}")
            return EvaluationResult(
                metrics=EvaluationMetrics(accuracy=0, fluency=0, consistency=0, terminology=0, completeness=0),
                suggestions="Failed to parse evaluation response",
                overall_score=0.0
            )




    def generate_reports(self, output_dir: str, report_prefix: str = "", original_pptx: str = None, translated_pptx: str = None):
        """Generate PDF, Excel, JSON reports, and Telemetry."""
        import os
        from .utils import generate_json_report as save_comparison_json, generate_excel_report, generate_pdf_report, save_model_mapping, generate_bilingual_pdf, save_token_usage, generate_visual_comparison_pdf
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Infer real model names from logs
        real_models = {
            "A": self.config.translator_config.model,
            "B": self.config.optimizer_config.model,
            "C": self.config.evaluator_config.model
        }
        
        for log in self.detailed_logs:
            tag = log.get("tag", "")
            if "Stage A" in tag:
                real_models["A"] = log.get("model", real_models["A"])
            elif "Stage C" in tag:
                real_models["B"] = log.get("model", real_models["B"])
            elif "Stage B" in tag: # Evaluator
                real_models["C"] = log.get("model", real_models["C"])
            
        # 1. Comparison JSON
        json_path = os.path.join(output_dir, f"{report_prefix}comparison.json")
        save_comparison_json(self.segments, json_path)
        
        # 2. Excel Report
        excel_path = os.path.join(output_dir, f"{report_prefix}evaluation_report.xlsx")
        generate_excel_report(self.segments, excel_path)
        
        # 3. PDF Report (Assessment)
        pdf_path = os.path.join(output_dir, f"{report_prefix}assessment_report.pdf")
        generate_pdf_report(self.segments, pdf_path, models=real_models)
        
        # 4. Bilingual PDF (Side-by-Side Text) - DISABLED per user request
        # bilingual_pdf_path = os.path.join(output_dir, f"{report_prefix}bilingual_text_comparison.pdf")
        # generate_bilingual_pdf(self.segments, bilingual_pdf_path)
        
        # 5. Visual Comparison PDF (Side-by-Side Slides)
        if original_pptx and translated_pptx:
            visual_pdf_path = os.path.join(output_dir, f"{report_prefix}comparison.pdf")
            generate_visual_comparison_pdf(original_pptx, translated_pptx, visual_pdf_path)
        
        # 5. Token Usage & Telemetry
        telemetry_path = os.path.join(output_dir, "token_usage.json")
        
        save_token_usage(self.detailed_logs, real_models, telemetry_path)
        
        # 6. Model Mapping
        save_model_mapping(output_dir, real_models)

